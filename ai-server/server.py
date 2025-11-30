# -*- coding: utf-8 -*-
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import httpx
from datetime import date
import uuid
from pptx import Presentation
import re
import json

# ======================================
# FastAPI 기본 설정
# ======================================
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ======================================
# 전역 버퍼
# ======================================
AI_PLAN_BUFFER = {}
# 구조:
# session_id → {
#     "plan_text": "...",
#     "todos": [...]
# }

# ======================================
# Ollama 설정
# ======================================
OLLAMA_URL = "http://127.0.0.1:11434/api/chat"
MODEL_NAME = "llama3.1:8b"


# ======================================
# 유틸 함수
# ======================================
def generate_session_id():
    return str(uuid.uuid4())


def detect_num_days(text: str, default=7) -> int:
    m = re.search(r"(\d+)\s*일", text)
    if not m:
        m = re.search(r"(\d+)\s*(day|days)", text, re.IGNORECASE)
    if m:
        try:
            return max(1, min(30, int(m.group(1))))
        except:
            return default
    return default


def extract_text_from_pptx(file_obj):
    try:
        prs = Presentation(file_obj)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print("⚠ PPT parsing error:", e)
        return ""


async def call_ollama(prompt: str, user_msg: str):
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {"role": "system", "content": prompt},
            {"role": "user", "content": user_msg},
        ],
        "stream": False,
        "num_predict": 8096,
    }

    async with httpx.AsyncClient(timeout=360) as client:
        res = await client.post(OLLAMA_URL, json=payload)
        res.raise_for_status()
        return res.json()["message"]["content"]


# ======================================
# 프롬프트 1: 텍스트용
# ======================================
def build_text_plan_prompt(num_days, today, ppt_text=None):

    ppt_block = f"\n[PPT 참고]\n{ppt_text[:]}\n" if ppt_text else ""

    return f"""
            당신은 한국어 학습 계획을 전문적으로 작성하는 AI입니다.

            [요구사항]
            - 총 {num_days}일 동안의 학습 계획을 작성합니다.
            - 하루별로 다음을 포함하세요:
            • 그날의 학습 목표 1개
            • 해야 할 구체적인 학습 작업 2~4개
            • 짧은 복습 또는 실습 항목 1개 포함
            - 모든 문장은 자연스럽고 실제 사람이 쓴 것처럼 구성합니다.
            - 전체 출력은 ‘한국어 설명 텍스트 ONLY’로 작성합니다.
            - JSON, 코드블록, 목록 마크다운 금지.
            - 날짜 정보(Today: {today})는 참고만 하세요.

            {ppt_block}

            한국어 학습 계획 작성을 시작하세요.
            """


# ======================================
# 프롬프트 2: JSON 생성용
# ======================================
def build_json_prompt(num_days, today, ppt_text=None):

    ppt_block = f"\n[PPT 참고]\n{ppt_text[:]}\n" if ppt_text else ""

    return f"""
            당신은 한국어 학습 계획을 JSON 형식으로 변환하는 역할만 수행합니다.
            아래 조건을 반드시 지키세요.

            [출력 형식(엄격히 준수)]
            JSON 배열 ONLY 출력:
            [
              {{
                "day": 1,
                "todos": [
                  {{
                    "title": "string",
                    "content": "string",
                    "status_id": "NOT_STARTED",
                    "end_time": "{today} 23:59",
                    "accumulated_time": 0
                  }},
                  ...
                ]
              }},
              {{
                "day": 2,
                "todos": [...same format...]
              }},
              ...
            ]

            [규칙]
            - JSON 외 텍스트 절대 출력 금지.
            - 코드블록, 해설, 불필요한 문장 추가 금지.
            - day는 1부터 {num_days}까지 순차 증가.
            - 각 day의 todos는 2~4개.
            - title, content는 한국어로 명확한 작업 설명.
            - PPT 참고 내용은 content에 자연스럽게 반영할 것.
            - JSON 문법 오류(쉼표/따옴표 누락 등) 절대 금지.

            {ppt_block}

            이제 위 형식대로 JSON만 출력하세요.
            """

# ======================================
# 프롬프트 3: 검색/질문 전용 챗봇
# ======================================
def build_search_prompt():
    return """
    당신은 한국어로만 답변하는 친절한 검색·지식 챗봇입니다.

    [규칙]
    - 사용자가 묻는 질문에 대해 간단하고 명확하게 답변하세요.
    - 공부/프로그래밍/일반 지식/상식 모두 답변 가능.
    - JSON, 코드블록, 목록 등은 요청할 때만 사용.
    - 학습 계획을 생성하지 말 것.
    - 투두 리스트 생성 금지.
    - 자연스럽고 부드러운 한국어 문장으로 답변.
    """

# ======================================
# (1) 텍스트 기반 플랜 생성
# ======================================
class ChatRequest(BaseModel):
    message: str


@app.post("/chat")
async def chat(req: ChatRequest):
    today = date.today()
    num_days = detect_num_days(req.message)

    # A) 텍스트 플랜 생성
    plan_text = await call_ollama(build_text_plan_prompt(num_days, today), req.message)

    # B) JSON 생성
    json_raw = await call_ollama(build_json_prompt(num_days, today), req.message)

    # JSON 파싱
    try:
        todos = json.loads(json_raw)
    except:
        todos = []

    session_id = generate_session_id()

    AI_PLAN_BUFFER[session_id] = {
        "plan_text": plan_text,
        "todos": todos,
    }

    return {
        "session_id": session_id,
        "answer": plan_text,
        "todo_count": len(todos),
        "todos": todos,
    }


# ======================================
# (2) PPT + 텍스트 기반 생성
# ======================================
@app.post("/chat-with-file")
async def chat_with_file(message: str = Form(...), file: UploadFile = File(...)):
    ppt_text = extract_text_from_pptx(file.file)
    today = date.today()
    num_days = detect_num_days(message)

    # A) 플랜 텍스트 생성
    plan_text = await call_ollama(build_text_plan_prompt(num_days, today, ppt_text), message)

    # B) JSON 생성
    json_raw = await call_ollama(build_json_prompt(num_days, today, ppt_text), message)

    try:
        todos = json.loads(json_raw)
    except:
        todos = []

    session_id = generate_session_id()
    AI_PLAN_BUFFER[session_id] = {
        "plan_text": plan_text,
        "todos": todos,
    }

    return {
        "success": True,
        "session_id": session_id,
        "answer": plan_text,
        "todo_count": len(todos),
        "todos": todos,
    }


# ======================================
# (3) Node가 JSON 요청할 때
# ======================================
@app.get("/get-json")
async def get_json(session_id: str):
    if session_id not in AI_PLAN_BUFFER:
        return {"success": False, "error": "Session not found"}

    return {
        "success": True,
        "todos": AI_PLAN_BUFFER[session_id]["todos"],
    }


# ======================================
# (4) 일반 지식 챗봇 (검색용)
# ======================================
class SearchChatRequest(BaseModel):
    message: str

@app.post("/ask")
async def ask(req: SearchChatRequest):
    prompt = build_search_prompt()
    answer = await call_ollama(prompt, req.message)

    return {
        "success": True,
        "answer": answer.strip()
    }


# ======================================
# (4) 성공 저장 후 버퍼 삭제
# ======================================
@app.post("/save-plan")
async def save_plan(session_id: str = Form(...)):
    if session_id in AI_PLAN_BUFFER:
        del AI_PLAN_BUFFER[session_id]
        return {"success": True}
    else:
        return {"success": False, "error": "Session not found"}


@app.get("/")
def root():
    return {"message": "Unified AI Server Running"}
